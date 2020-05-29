
import xlrd
import textract
import PyPDF2 
from pptx import Presentation
import glob
import csv
import os 
import shutil

def Search_xlsx(value_,filename):

	#FOR OPENING THE WORKBOOK
	loc=('path of file')
	wb=xlrd.open_workbook(filename)
	sheet=wb.sheet_by_index(0)
	sheet.cell_value(0,0)


	#print(sheet.nrows)
	#print(sheet.ncols)

	rows_=sheet.nrows	#no. of rows
	cols_=sheet.ncols	#no. of columns

	# print(sheet.cell_value(0,0))
	# print(sheet.cell_value(0,1))

	for el in range(cols_):
		for elm in range(rows_):
			if value_.lower()==str(sheet.cell_value(elm,el)).lower():
				return True			
			#print(sheet.cell_value(elm,el))

def Search_doc(value_,filename):
	text=textract.process(filename)
	if value_.lower() in str(text).lower():
		return True
	else:
		return False

def Search_pdf(value_,filename):
	# creating a pdf file object 
	pdfFileObj = open(filename, 'rb') 
	  
	# creating a pdf reader object 
	pdfReader = PyPDF2.PdfFileReader(pdfFileObj) 
	  
	# printing number of pages in pdf file 
	#print(pdfReader.numPages) 
	  
	# creating a page object 
	pageObj = pdfReader.getPage(0) 
	  
	# extracting text from page 
	#print(pageObj.extractText())

	string_=pageObj.extractText()

	#print (value_)

	#converted the string in lower case and
	#also converted the pdf text to lower case so that
	#we case get case insensitive fucntionality
	if value_.lower() in str(string_.lower()):
		return True
	else:
	 	return False	
	  
	# closing the pdf file object 
	pdfFileObj.close()

def Search_PPT(value_,filename):
	prs=Presentation(filename)
	for slide in prs.slides:
		for shape in slide.shapes:
			if hasattr(shape,"text"):
				if value_.lower() in shape.text.lower():
					return True
				else:
					return False	


def Search_CSV(value_,filename):
	text_string=[]
	with open(filename, newline='') as csvfile:
		csv_reader=csv.reader(csvfile,delimiter=' ',quotechar='|')
		for row in csv_reader:
			for el in row:
				if value_.lower() in el.lower():
					return True
				else:
					return False			
		# if value_.lower() in text_string.lower():
		# 	return True
		# else:
		# 	return False		

def filename_Search(value_):
	dir_path = os.path.dirname(os.path.realpath(__file__))
	print("Using glob.glob()") 
	files = glob.glob('*{filename}*.*'.format(filename=value_),  
	                   recursive = True) 
	for file in files: 
	    print(file)
	    shutil.copy(file,dir_path+'\\'+"final_sorted"+"\\")
	    #os.system('copy '+file+" "+dir_path+'\\'+"final_sorted"+"\\")
	

def Search_text(value_,filename):
	print("dffffffffffffffffffffffffffffff",filename)
	f=open(filename,'r')
	text_content=f.read()
	if value_.lower() in str(text_content).lower():
		print("################################3")
		return True
	else:
		return False
	f.close()		

def getting_allfiles(value_):
	# This is to get the directory that the program  
	# is currently running in. 
	dir_path = os.path.dirname(os.path.realpath(__file__))
	#print("\n\n DIR path",dir_path)   
	for root, dirs, files in os.walk('.'):
		if root==".\\final_sorted":
			continue
		for file in files:
			try:
				path2=os.getcwd()
				#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
				#path_=path2+'\\'+file
				path_=os.path.join(root, file)
				#if file.endswith('.docx') or file.endswith('.xlsx') or file.endswith('.pdf') or file.endswith('.pptx') or file.endswith('.csv'):
					#print(root+'\\'+str(file))
					
					#print("\n\n#################\n#################")

				if file.endswith('docx'):	
					truth=Search_doc(value_,path_)
					if truth==True:
						#print("path-----------",path_)
						print('\n Key word found in doc file File NAME IS:',file)
						shutil.copy(path_,dir_path+'\\'+"final_sorted"+"\\")
						#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
						#os.system('copy '+path_+" "+dir_path+'\\'+"final_sorted"+"\\")
				
				if file.endswith('.pdf'):
					truth=Search_pdf(value_,path_)
					if truth==True:
						print("\nkey word found in pdf file ,File NAME IS: ",file)
						#print("PATH TO PDF_>",path_)
						shutil.copy(path_,dir_path+'\\'+"final_sorted"+"\\")
						#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
						#os.system('copy '+path_+" "+dir_path+'\\'+"final_sorted"+"\\")
				if file.endswith('.xlsx'):
					truth=Search_xlsx(value_,path_)
					if truth==True:
						print("\nkey word found in excel file,File NAME IS: ",file)
						shutil.copy(path_,dir_path+'\\'+"final_sorted"+"\\")
						#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
						#os.system('copy '+path_+" "+dir_path+'\\'+"final_sorted"+"\\")																		
				if file.endswith('.pptx'):
					truth=Search_PPT(value_,path_)
					if truth==True:
						print("\nkey word found in PPT file,File NAME IS: ",file)
						shutil.copy(path_,dir_path+'\\'+"final_sorted"+"\\")
						#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
						#os.system('copy '+path_+" "+dir_path+'\\'+"final_sorted"+"\\")
				if file.endswith('.txt'):
					truth=Search_text(value_,path_)
					if truth==True:
						print("\nKey word found in text file, File name is :",file)
						shutil.copy(path_,dir_path+"\\"+"final_sorted"+"\\")		
				if file.endswith('.csv'):
					truth=Search_CSV(value_,path_)
					if truth==True:
						print("\nkey word found in csv file ",file)
						shutil.copy(path_,dir_path+'\\'+"final_sorted"+"\\")
						#print("\nroot",root,"\ndirs",dirs,"\nfiles",files)
						#os.system('copy '+path_+" "+dir_path+'\\'+"final_sorted"+"\\")

					#since permission is denied when i give full path from c direcorty
					#shutil.copyfile(root+'\\'+str(file),dir_path+"\\"+"final_sorted"+"\\")
			except Exception as e:
				print(file,e)



        # # change the extension from '.mp3' to  
        # # the one of your choice. 
        # if file.endswith('.mp3'): 
        #     print root+'/'+str(file)
flag=1
while flag==1:
	if os.path.isdir("final_sorted")==False:
		os.system("mkdir final_sorted")	
	print("-------------------------------------")
	inp1=int(input("\na.Enter 1 for search in content of files\nb.Enter 2 for content plus search in file names\nc.Enter 0 to exit\n:"))	

	if inp1==1:
		input_string=str(input("Enter the string to be searched"))
		getting_allfiles(input_string)
	if inp1==2:
		input_string=str(input("Enter the string to be searched"))
		getting_allfiles(input_string)
		filename_Search(input_string)
	if inp1==0:
		break	
# var1=Search_xlsx("vikas","bhai ka exel.xlsx")
# var2=Search_doc("vikas",'mera.docx')
# var3=Search_pdf("vikas",'Offer Letter Vikas.pdf')
# var4=Search_PPT('cyber',"Fundamentals of Cyber Forensics.pptx")
# var5=Search_CSV("cyber",'bhai ka exel.csv')
#print("Search_xlsx/excel:",var1,"\nSearch_doc/word:",var2,"\nSearch_pdf:",var3,"\nSearch_PPT:",var4,"\nSearch_CSV:",var5)	

#getting_allfiles('vikas')
#var9=filename_Search('funda')



	# doc_files=glob.glob('*.docx',recursive=True)
	# excel_files=glob.glob('*.xlsx',recursive=True)
	# pdf_files=glob.glob('*.pdf',recursive=True)
	# csv_files=glob.glob('*.csv',recursive=True)
	# ppt_files=glob.glob('*.pptx',recursive=True)
	# for file in doc_files:
	# 	print(file)
	# print(doc_files,excel_files,pdf_files,csv_files,ppt_files)	